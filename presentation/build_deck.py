#!/usr/bin/env python3
"""
MR-Fit — Graduation Defense Deck Generator
==========================================
Builds a premium, startup-launch-style PowerPoint (.pptx) from the project's
screenshots folder.

USAGE
-----
    pip install python-pptx
    python presentation/build_deck.py

OUTPUT
------
    presentation/MR-Fit_Graduation_Defense.pptx

NOTES
-----
* Images are pulled from ../screenshots relative to this file.
* Morph transitions + Slide Zoom must be enabled in PowerPoint (see README in
  this folder). python-pptx cannot write Morph/Zoom XML, so the deck is built so
  that turning Morph ON (Transitions > Morph, Apply To All) gives the Prezi-style
  camera feel because objects keep consistent positions/sizes across slides.
* 6 presenter sections are colour-tabbed and labelled; speaker notes per slide.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Brand system
# ----------------------------------------------------------------------------
BG        = RGBColor(0x0B, 0x0F, 0x14)   # near-black canvas
BG2       = RGBColor(0x11, 0x18, 0x20)   # panel
GREEN     = RGBColor(0x22, 0xE3, 0x7D)   # primary accent (energy)
CYAN      = RGBColor(0x3B, 0xC9, 0xFF)   # secondary accent
WHITE     = RGBColor(0xE8, 0xED, 0xF2)
MUTED     = RGBColor(0x8A, 0x97, 0xA6)
RED       = RGBColor(0xFF, 0x5C, 0x5C)

HEAD_FONT = "Sora"            # falls back gracefully if not installed
ALT_HEAD  = "Space Grotesk"
BODY_FONT = "Inter"
MONO_FONT = "JetBrains Mono"

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "..", "screenshots")

def shot(name):
    return os.path.join(SHOTS, name)

# Presenter colour tabs
PRESENTERS = {
    1: ("Presenter 1", GREEN),
    2: ("Presenter 2", CYAN),
    3: ("Presenter 3", RGBColor(0xB5, 0x8CFF & 0xFF, 0xFF)),
    4: ("Presenter 4", RGBColor(0xFF, 0xB3, 0x4D)),
    5: ("Presenter 5", RGBColor(0xFF, 0x7A, 0xB3)),
    6: ("Presenter 6", RGBColor(0x6E, 0xE7, 0xB7)),
}

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------

def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def rect(slide, x, y, w, h, color, line=None, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1.25)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, text, size, color=WHITE, bold=False,
            font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            italic=False, spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = font; r.font.color.rgb = color
    return tb


def accent_bar(slide, color=GREEN, x=Inches(0.6), y=Inches(1.55), w=Inches(1.4)):
    rect(slide, x, y, w, Pt(5), color)


def presenter_tab(slide, num, section_label):
    name, color = PRESENTERS[num]
    # top-right pill
    pill = rect(slide, SW - Inches(3.4), Inches(0.35), Inches(3.0), Inches(0.5),
                BG2, rounded=True)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, SW - Inches(3.25),
                                 Inches(0.47), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    dot.shadow.inherit = False
    textbox(slide, SW - Inches(2.9), Inches(0.34), Inches(2.6), Inches(0.5),
            f"{name.upper()}  •  {section_label}", 10.5, color, bold=True,
            font=MONO_FONT, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, page):
    textbox(slide, Inches(0.6), SH - Inches(0.5), Inches(4), Inches(0.4),
            "MR-Fit — AI Fitness Companion", 9, MUTED, font=MONO_FONT)
    textbox(slide, SW - Inches(1.4), SH - Inches(0.5), Inches(1), Inches(0.4),
            f"{page:02d}", 9, MUTED, font=MONO_FONT, align=PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_image_card(slide, path, x, y, w, h, border=GREEN):
    """Place an image inside a thin neon frame; scales to fit box."""
    frame = rect(slide, x - Pt(3), y - Pt(3), w + Pt(6), h + Pt(6), BG2,
                 line=border, rounded=True)
    try:
        slide.shapes.add_picture(path, x, y, width=w, height=h)
    except Exception as e:
        textbox(slide, x, y + h/2, w, Inches(0.5),
                f"[missing: {os.path.basename(path)}]", 12, MUTED,
                align=PP_ALIGN.CENTER)
    return frame

PAGE = 0

def pageno():
    global PAGE; PAGE += 1; return PAGE

# ============================================================================
# SLIDE 0 — COLD OPEN / TITLE  (Presenter 1)
# ============================================================================
s = add_slide(); fill_bg(s)
# faint orbital ring
ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.6), Inches(0.2),
                          Inches(6.1), Inches(6.1))
ring.fill.background(); ring.line.color.rgb = BG2; ring.line.width = Pt(1.5)
ring.shadow.inherit = False
try:
    s.shapes.add_picture(shot("logo_transparent.png"), Inches(5.16), Inches(1.45),
                         height=Inches(2.6))
except Exception:
    pass
textbox(s, Inches(1), Inches(4.25), Inches(11.33), Inches(1),
        "MR-FIT", 60, WHITE, bold=True, font=HEAD_FONT, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(5.35), Inches(11.33), Inches(0.6),
        "Your AI-Powered Fitness Companion", 22, GREEN, font=BODY_FONT,
        align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(6.5), Inches(11.33), Inches(0.5),
        "Graduation Project Defense  •  Team MR-Fit  •  2026", 12, MUTED,
        font=MONO_FONT, align=PP_ALIGN.CENTER)
notes(s, "PRESENTER 1 (Cold open, ~0:45). Hook: 'How many fitness apps have "
          "you downloaded — and deleted?' Pause, then reveal MR-Fit. Keep energy "
          "high; do not read the slide. Set up that fitness apps are everywhere "
          "but intelligent ones are not.")
footer(s, pageno())

# ============================================================================
# SLIDE 1 — THE PROBLEM  (Presenter 1)
# ============================================================================
s = add_slide(); fill_bg(s)
presenter_tab(s, 1, "THE PROBLEM")
textbox(s, Inches(0.6), Inches(0.7), Inches(10), Inches(0.9),
        "Fitness apps are everywhere.", 40, WHITE, bold=True, font=HEAD_FONT)
textbox(s, Inches(0.6), Inches(1.55), Inches(11), Inches(0.7),
        "Intelligent ones aren’t.", 30, GREEN, bold=True, font=HEAD_FONT)
cards = [
    ("GENERIC", "One template for millions of different bodies and goals."),
    ("STATIC", "Plans never adapt to how your week, sleep, or recovery actually went."),
    ("SILOED", "Workouts, nutrition, and wearable data never talk to each other."),
]
cx = Inches(0.6)
for i, (title, body) in enumerate(cards):
    x = Inches(0.6 + i * 4.15)
    rect(s, x, Inches(2.9), Inches(3.85), Inches(3.0), BG2, line=RED, rounded=True)
    textbox(s, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(0.6),
            title, 22, RED, bold=True, font=HEAD_FONT)
    textbox(s, x + Inches(0.3), Inches(4.0), Inches(3.3), Inches(1.8),
            body, 16, WHITE, font=BODY_FONT)
notes(s, "PRESENTER 1 (~1:30). Walk through the three pains slowly, one card at "
          "a time. Land the line: the problem isn't tracking — it's that nothing "
          "thinks. Transition: 'What if one platform could see everything and "
          "actually think? Here's [Presenter 2].'")
footer(s, pageno())

# ============================================================================
# SLIDE 2 — THE VISION  (Presenter 2)
# ============================================================================
s = add_slide(); fill_bg(s)
presenter_tab(s, 2, "THE VISION")
textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
        "One platform that sees everything — and adapts.", 34, WHITE, bold=True,
        font=HEAD_FONT)
accent_bar(s, CYAN, y=Inches(1.85))
textbox(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.0),
        "MR-Fit unifies your training, nutrition, and wearable data, then lets "
        "AI act on the full picture — generating plans, coaching you, and "
        "counting your reps automatically.", 20, MUTED, font=BODY_FONT)
for i, t in enumerate(["UNIFY", "UNDERSTAND", "ADAPT"]):
    x = Inches(1.4 + i * 3.7)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(3.7), Inches(2.6), Inches(2.6))
    o.fill.solid(); o.fill.fore_color.rgb = BG2
    o.line.color.rgb = CYAN if i == 1 else GREEN; o.line.width = Pt(2)
    o.shadow.inherit = False
    textbox(s, x, Inches(4.75), Inches(2.6), Inches(0.6), t, 18, WHITE,
            bold=True, font=HEAD_FONT, align=PP_ALIGN.CENTER)
notes(s, "PRESENTER 2 (~1:30). Introduce MR-Fit as the AI fitness companion. "
          "Core idea: unify the data, then let AI act on it. Keep it visionary "
          "and confident. Then move into the product tour.")
footer(s, pageno())

# ============================================================================
# FEATURE SLIDES  (Presenters 2 & 3)
# ============================================================================

def feature_slide(presenter, label, title, bullets, img, img_left=True,
                  accent=GREEN):
    s = add_slide(); fill_bg(s)
    presenter_tab(s, presenter, label)
    img_w, img_h = Inches(6.1), Inches(4.6)
    if img_left:
        ix, iy = Inches(0.7), Inches(1.9)
        tx = Inches(7.3)
    else:
        ix, iy = Inches(6.5), Inches(1.9)
        tx = Inches(0.7)
    add_image_card(s, shot(img), ix, iy, img_w, img_h, border=accent)
    textbox(s, tx, Inches(0.8), Inches(5.6), Inches(0.5),
            label, 13, accent, bold=True, font=MONO_FONT)
    textbox(s, tx, Inches(1.25), Inches(5.6), Inches(0.9),
            title, 30, WHITE, bold=True, font=HEAD_FONT)
    by = 2.5 if img_left else 2.5
    bullet_box = textbox(s, tx, Inches(2.6), Inches(5.6), Inches(4.0), "",
                         16, WHITE, font=BODY_FONT)
    tf = bullet_box.text_frame
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15; p.space_after = Pt(10)
        r = p.add_run(); r.text = "▸  " + b
        r.font.size = Pt(16); r.font.name = BODY_FONT; r.font.color.rgb = WHITE
    return s

s = feature_slide(2, "FEATURE — LANDING", "A launch that sells the vision",
    ["Modern marketing landing page", "Clear value proposition up front",
     "Sign-in / sign-up flow"], "landing_page.png", img_left=True, accent=GREEN)
notes(s, "PRESENTER 2. Show the polished entry point — this is a product, not a "
          "school project. Mention auth (NextAuth / Supabase JWT).")
footer(s, pageno())

s = feature_slide(2, "FEATURE — DASHBOARD", "Everything in one glance",
    ["Unified home for training + nutrition", "Daily macros vs. goals",
     "Quick access to AI tools"], "dashboard.png", img_left=False, accent=GREEN)
notes(s, "PRESENTER 2. The dashboard is where the 'unify' promise becomes real. "
          "Transition to Presenter 3 for the deep feature tour.")
footer(s, pageno())

s = feature_slide(3, "FEATURE — WORKOUTS", "Log every set, rep & weight",
    ["Create and log workouts in seconds", "Sets, reps, weights and notes",
     "AI-generated plans flagged as source: ai"], "active_workout.png",
    img_left=True, accent=CYAN)
notes(s, "PRESENTER 3 (~start). Walk through logging UX. Note plans can be AI- or "
          "user-created.")
footer(s, pageno())

s = feature_slide(3, "FEATURE — NUTRITION", "Eat with intent",
    ["Search a large food database", "Track daily macros against goals",
     "Nutrition and training in one place"], "nutrition logging.png",
    img_left=False, accent=CYAN)
notes(s, "PRESENTER 3. Emphasise nutrition + training living together — the data "
          "silo problem solved.")
footer(s, pageno())

s = feature_slide(3, "FEATURE — PROGRESS", "See the trend, not just the day",
    ["Weight trend chart", "Workout frequency heatmap",
     "Personal records table"], "progress tracker.png", img_left=True, accent=CYAN)
notes(s, "PRESENTER 3. Progress visualisation drives motivation. Transition: "
          "'All of this feeds the part that makes us different — the AI brain. "
          "Here's [Presenter 4].'")
footer(s, pageno())

# ============================================================================
# AI SECTION  (Presenter 4)
# ============================================================================
s = add_slide(); fill_bg(s)
presenter_tab(s, 4, "THE AI BRAIN")
textbox(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.8),
        "The AI Brain", 40, WHITE, bold=True, font=HEAD_FONT)
textbox(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.6),
        "Three AI systems working on your real data.", 18, CYAN, font=BODY_FONT)
ai_cards = [
    ("RAG WORKOUT GENERATION", GREEN,
     "800+ exercises embedded as 1536-dim vectors (pgvector). On 'Generate Plan', "
     "a similarity search picks the top candidates, then GPT-4o returns a "
     "structured JSON plan tailored to your goal and level."),
    ("AI COACH", CYAN,
     "A chat coach that reads live wearable context — sleep, HRV, resting heart "
     "rate — to give advice grounded in how your body is actually doing."),
    ("SMART TRACKER (EDGE AI)", RGBColor(0xFF,0xB3,0x4D),
     "IMU sensor data streams into a trained Random Forest model (FastAPI) that "
     "classifies the exercise and counts reps automatically."),
]
for i, (t, c, b) in enumerate(ai_cards):
    x = Inches(0.6 + i * 4.15)
    rect(s, x, Inches(2.3), Inches(3.85), Inches(4.4), BG2, line=c, rounded=True)
    textbox(s, x + Inches(0.3), Inches(2.55), Inches(3.3), Inches(1.0),
            t, 16, c, bold=True, font=HEAD_FONT)
    textbox(s, x + Inches(0.3), Inches(3.7), Inches(3.3), Inches(2.8),
            b, 14, WHITE, font=BODY_FONT)
notes(s, "PRESENTER 4 (~3:30). The crown jewel. Explain WHY RAG (grounds the LLM "
          "in our real exercise catalog, cuts hallucination). Be ready for: how "
          "does the Random Forest rep-counter work, what accuracy. Keep it "
          "concrete.")
footer(s, pageno())

# AI feature screenshots (form analysis + recovery + coach)
s = feature_slide(4, "AI — FORM ANALYSIS", "AI form analysis",
    ["Computer-vision / sensor-driven feedback", "Flags technique issues",
     "Helps prevent injury"], "ai form analysis.png", img_left=True, accent=CYAN)
notes(s, "PRESENTER 4. Show form analysis output. Tie back to Smart Tracker data.")
footer(s, pageno())

s = feature_slide(4, "AI — RECOVERY ENGINE", "AI recovery engine",
    ["Reads wearable recovery signals", "Recommends rest vs. push",
     "Adapts the plan to your body"], "ai reckovery engine.png", img_left=False,
    accent=RGBColor(0xFF,0xB3,0x4D))
notes(s, "PRESENTER 4. The recovery engine closes the loop — data in, smarter "
          "plan out. Transition to Presenter 5 for architecture.")
footer(s, pageno())

s = feature_slide(4, "AI — COACH", "Chat that knows your data",
    ["Conversational fitness guidance", "Pulls live wearable context",
     "Personalized, not generic"], "ai_coach.png", img_left=True, accent=CYAN)
notes(s, "PRESENTER 4. Demo-style narration of the coach. Then hand to Presenter 5.")
footer(s, pageno())

# ============================================================================
# ARCHITECTURE  (Presenter 5)
# ============================================================================
s = add_slide(); fill_bg(s)
presenter_tab(s, 5, "ARCHITECTURE")
textbox(s, Inches(0.6), Inches(0.55), Inches(12), Inches(0.8),
        "How it all fits together", 34, WHITE, bold=True, font=HEAD_FONT)

layers = [
    ("CLIENT", "Next.js 14 App Router — Server Components, Tailwind UI", GREEN),
    ("EDGE / BFF", "Next.js Route Handlers — validate & proxy requests", CYAN),
    ("DATA — SUPABASE", "Postgres 15 • Row-Level Security • Realtime • pgvector • Storage", GREEN),
    ("AUTOMATION — n8n", "workout-plan-generator • wearable-sync (cron) • progress-report", CYAN),
    ("AI / EXTERNAL", "OpenAI GPT-4o + embeddings • Garmin / Fitbit / Apple Health", RGBColor(0xFF,0xB3,0x4D)),
]
y = 1.6
for name, desc, c in layers:
    rect(s, Inches(0.8), Inches(y), Inches(11.7), Inches(0.92), BG2, line=c, rounded=True)
    textbox(s, Inches(1.1), Inches(y+0.12), Inches(3.3), Inches(0.7),
            name, 16, c, bold=True, font=MONO_FONT, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.5), Inches(y+0.12), Inches(7.7), Inches(0.7),
            desc, 14, WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
    y += 1.07
notes(s, "PRESENTER 5 (~3:00). Build the stack top-to-bottom. Stress security: "
          "Row-Level Security isolates user data at the DB layer; service-role key "
          "never reaches the browser. Be ready for 'why Supabase + n8n'.")
footer(s, pageno())

# ============================================================================
# DEMO  (Presenter 6)
# ============================================================================
s = add_slide(); fill_bg(s)
presenter_tab(s, 6, "LIVE DEMO")
rect(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.2), BG2, line=GREEN, rounded=True)
textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.0),
        "LIVE DEMO", 48, GREEN, bold=True, font=HEAD_FONT, align=PP_ALIGN.CENTER)
textbox(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.6),
        "Sign in  →  Generate AI plan  →  Log a set  →  Watch the chart update live  "
        "→  Ask the AI Coach  →  Open the wearables dashboard", 18, WHITE,
        font=BODY_FONT, align=PP_ALIGN.CENTER)
notes(s, "PRESENTER 6 drives; Presenters 2 & 4 narrate features/AI. HAVE A BACKUP "
          "VIDEO and pre-seed demo_seed.sql (17 workouts, 30 nutrition entries, "
          "15 weight logs) in case of network issues.")
footer(s, pageno())

# ============================================================================
# ROADMAP + RESULTS  (Presenter 6)
# ============================================================================
s = add_slide(); fill_bg(s)
presenter_tab(s, 6, "RESULTS & ROADMAP")
textbox(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.8),
        "What we built — and what’s next", 32, WHITE, bold=True, font=HEAD_FONT)
col = [
    ("DELIVERED", GREEN, ["Full-stack AI fitness platform",
        "RAG workout generation", "Realtime tracking + charts",
        "Edge-AI rep counting", "Wearables integration"]),
    ("ROADMAP", CYAN, ["Native mobile app",
        "More wearable providers", "Social / community challenges",
        "On-device model for offline", "Coach voice mode"]),
]
for i, (t, c, items) in enumerate(col):
    x = Inches(0.7 + i * 6.2)
    rect(s, x, Inches(1.8), Inches(5.8), Inches(4.7), BG2, line=c, rounded=True)
    textbox(s, x + Inches(0.4), Inches(2.05), Inches(5), Inches(0.6),
            t, 20, c, bold=True, font=HEAD_FONT)
    bb = textbox(s, x + Inches(0.4), Inches(2.85), Inches(5), Inches(3.4), "",
                 16, WHITE, font=BODY_FONT)
    tf = bb.text_frame
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r = p.add_run(); r.text = "▸  " + it
        r.font.size = Pt(16); r.font.name = BODY_FONT; r.font.color.rgb = WHITE
notes(s, "PRESENTER 6. Summarise impact and learning. Be ready for: 'is it "
          "production-ready, hardest problem solved, what's next'.")
footer(s, pageno())

# ============================================================================
# CLOSE  (Presenter 6)
# ============================================================================
s = add_slide(); fill_bg(s)
try:
    s.shapes.add_picture(shot("logo_transparent.png"), Inches(5.66), Inches(1.6),
                         height=Inches(2.1))
except Exception:
    pass
textbox(s, Inches(1), Inches(3.9), Inches(11.33), Inches(0.9),
        "Thank you.", 44, WHITE, bold=True, font=HEAD_FONT, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.9), Inches(11.33), Inches(0.6),
        "MR-Fit — a fitness companion that sees your whole picture and adapts.",
        18, GREEN, font=BODY_FONT, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(5.9), Inches(11.33), Inches(0.5),
        "Team MR-Fit  •  Questions welcome", 13, MUTED, font=MONO_FONT,
        align=PP_ALIGN.CENTER)
notes(s, "PRESENTER 6 closing line: 'MR-Fit isn't just another tracker — it's a "
          "companion that sees your whole picture and adapts. We're Team MR-Fit. "
          "Thank you — we'd love your questions.'")
footer(s, pageno())

# ----------------------------------------------------------------------------
# Save
# ----------------------------------------------------------------------------
out = os.path.join(HERE, "MR-Fit_Graduation_Defense.pptx")
prs.save(out)
print(f"Saved: {out}")
print("Next: open in PowerPoint, select all slides, Transitions > Morph > Apply To All.")

